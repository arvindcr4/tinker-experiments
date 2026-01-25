#!/usr/bin/env python3
"""Create PowerPoint presentation for Tinker RL demo."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)

    # Bullets
    bullet_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.5), Inches(5.5))
    tf = bullet_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.space_after = Pt(12)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    return slide

def add_table_slide(prs, title, headers, rows):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)

    # Table
    cols = len(headers)
    table = slide.shapes.add_table(len(rows) + 1, cols, Inches(0.5), Inches(1.3), Inches(9), Inches(0.5 * (len(rows) + 1))).table

    # Headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)

    # Rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(12)

    return slide

def add_code_slide(prs, title, code):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True

    # Code box with background
    code_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(1.2), Inches(9.4), Inches(5.3))
    code_shape.fill.solid()
    code_shape.fill.fore_color.rgb = RGBColor(0x28, 0x28, 0x28)
    code_shape.line.fill.background()

    # Code text
    code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5))
    tf = code_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(14)
    p.font.name = "Courier New"
    p.font.color.rgb = RGBColor(0x00, 0xff, 0x00)

    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs,
        "Reinforcement Learning for LLMs with Tinker",
        "Training Language Models with Human Feedback\n\nPES LLM Research Project")

    # Slide 2: Overview
    add_content_slide(prs, "What is Tinker?", [
        "Cloud-based RL training platform for LLMs",
        "Supports LoRA fine-tuning with reinforcement learning",
        "Provides pre-built recipes: Math RL, Preference Learning, Distillation",
        "No GPU required locally - training runs in the cloud",
        "API-based access with Python SDK"
    ])

    # Slide 3: Training Methods
    add_content_slide(prs, "Training Methods Explored", [
        "Supervised Fine-Tuning (SFT): Train on human-written examples",
        "Reinforcement Learning (RL): Learn from reward signals",
        "Preference Learning: Learn from pairwise comparisons",
        "Distillation: Transfer knowledge from teacher to student",
        "DPO: Direct Preference Optimization"
    ])

    # Slide 4: Experiments Table
    add_table_slide(prs, "Experiments Run",
        ["Recipe", "Task", "Model", "Key Result"],
        [
            ["Math RL", "Arithmetic", "Llama-3.2-1B", "0.67 → 1.0 reward"],
            ["Chat SL", "Conversation", "Llama-3.2-1B", "NLL decreasing"],
            ["Preference", "Shorter responses", "Qwen-0.6B", "Learning conciseness"],
            ["Distillation", "Reasoning transfer", "Llama-3.2-1B", "KL minimization"],
            ["GSM8K RL", "Word problems", "Llama-3.2-1B", "Multi-step reasoning"],
        ])

    # Slide 5: Math RL Results
    add_content_slide(prs, "Math RL: Arithmetic Results", [
        "Task: Train model to add two-digit numbers correctly",
        "Starting accuracy: 69.5% | Final accuracy: 100%",
        "Reward improved: 0.676 → 1.0 in ~20 steps",
        "Training time: ~2 minutes total",
        "Key insight: RL quickly masters simple arithmetic"
    ])

    # Slide 6: Code Example
    add_code_slide(prs, "Running Math RL Training",
"""python -m tinker_cookbook.recipes.math_rl.train \\
    model_name="meta-llama/Llama-3.2-1B" \\
    env=arithmetic \\
    group_size=4 \\
    groups_per_batch=100 \\
    learning_rate=1e-4

# Output:
# Step 0:  reward=0.676, correct=69.5%
# Step 10: reward=0.998, correct=99.8%
# Step 20: reward=1.000, correct=100%""")

    # Slide 7: Preference Learning
    add_content_slide(prs, "Preference Learning: Shorter Responses", [
        "Goal: Train model to generate concise responses",
        "Method: Pairwise comparison - shorter response wins",
        "Reward structure: win_minus_loss + format penalty",
        "Observation: Model learns to end responses properly",
        "Trade-off: Balancing brevity with informativeness"
    ])

    # Slide 8: Distillation
    add_content_slide(prs, "Knowledge Distillation", [
        "Off-Policy: Train on pre-collected teacher traces (SFT)",
        "On-Policy: Student generates, minimizes KL to teacher",
        "Dataset: OpenThoughts3 (1.2M reasoning traces)",
        "Metric: Teacher KL divergence decreases over training",
        "Benefit: Smaller student learns from larger teacher"
    ])

    # Slide 9: Technical Architecture
    add_content_slide(prs, "Technical Architecture", [
        "API Key authentication with Tinker cloud",
        "LoRA adapters for efficient fine-tuning (rank 32-128)",
        "Batch sampling with configurable group sizes",
        "Metrics logged to JSONL files and optional W&B",
        "Checkpoint saving for resumable training"
    ])

    # Slide 10: Jupyter Notebooks
    add_content_slide(prs, "Jupyter Notebooks Created", [
        "01_math_rl_arithmetic.ipynb - Basic RL training",
        "02_chat_sl_sft.ipynb - Supervised fine-tuning",
        "03_preference_shorter.ipynb - Preference learning",
        "04_distillation_off_policy.ipynb - SFT distillation",
        "05_distillation_on_policy.ipynb - KL distillation",
        "06_math_rl_gsm8k.ipynb - Word problem solving"
    ])

    # Slide 11: Key Findings
    add_content_slide(prs, "Key Findings", [
        "Simple tasks (arithmetic): RL achieves 100% quickly",
        "Complex tasks (GSM8K): Requires more training, larger models",
        "Preference learning: Effective for stylistic control",
        "Distillation: Efficient transfer from teacher models",
        "Cloud training: Accessible without local GPU resources"
    ])

    # Slide 12: Future Work
    add_content_slide(prs, "Future Directions", [
        "Code generation with sandbox execution (DeepCoder)",
        "Multi-agent RL training (games, debates)",
        "RLHF with human preference data",
        "Scaling to larger models (8B, 32B parameters)",
        "Custom reward functions for domain-specific tasks"
    ])

    # Slide 13: Conclusion
    add_title_slide(prs,
        "Thank You!",
        "Questions?\n\nPES LLM Research Project\nTinker Cookbook Experiments")

    # Save
    prs.save('/home/ubuntu/glass/Tinker_RL_Demo.pptx')
    print("Presentation saved to /home/ubuntu/glass/Tinker_RL_Demo.pptx")

if __name__ == "__main__":
    main()
